import requests
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import sys
import os

OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "NeuralNexusLab/HacKing:latest"

# Professional color scheme
COLORS = {
    'primary': RGBColor(26, 115, 232),
    'secondary': RGBColor(234, 67, 53),
    'accent': RGBColor(251, 188, 5),
    'success': RGBColor(52, 168, 83),
    'dark': RGBColor(32, 33, 36),
    'light': RGBColor(241, 243, 244),
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(60, 64, 67),
}


def emit_progress(step, message, data=None):
    """Send progress update to stdout as JSON line."""
    payload = {"step": step, "message": message}
    if data:
        payload["data"] = data
    print(json.dumps(payload), flush=True)


def call_ollama(prompt):
    """Send a prompt to Ollama and return the cleaned response text."""
    response = requests.post(
        OLLAMA_URL,
        json={"model": MODEL_NAME, "prompt": prompt, "stream": False},
        timeout=120
    )
    response.raise_for_status()
    data = response.json()
    raw = data.get("response", "")
    # Strip <think>...</think> blocks
    raw = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL)
    return raw.strip()


# ============================================================
# STEP 1: Understand the user's prompt and extract a clear topic
# ============================================================
def step1_extract_topic(user_prompt):
    """Ask the LLM to interpret the user's prompt and return a clear topic."""
    emit_progress(1, "Understanding your prompt...")

    prompt = f"""The user wants to create a presentation. Their input is:
"{user_prompt}"

Your job: figure out the exact presentation topic from their input.
Return ONLY a JSON object with these fields, nothing else:
{{"topic": "<clear topic title>", "subtitle": "<short subtitle>", "num_slides": <number between 5 and 12>}}

Example:
{{"topic": "Introduction to Artificial Intelligence", "subtitle": "A beginner-friendly overview", "num_slides": 8}}

Return ONLY the JSON:"""

    raw = call_ollama(prompt)

    # Clean and parse
    raw = re.sub(r"```json\s*", "", raw)
    raw = re.sub(r"```\s*", "", raw)
    match = re.search(r"\{.*\}", raw, re.DOTALL)

    if match:
        try:
            result = json.loads(match.group())
            topic = result.get("topic", user_prompt)
            subtitle = result.get("subtitle", "")
            num_slides = min(max(int(result.get("num_slides", 7)), 3), 15)
            emit_progress(1, f"Topic identified: {topic}", {
                "topic": topic, "subtitle": subtitle, "num_slides": num_slides
            })
            return topic, subtitle, num_slides
        except (json.JSONDecodeError, ValueError):
            pass

    # Fallback: use the prompt as-is
    emit_progress(1, f"Topic identified: {user_prompt}", {
        "topic": user_prompt, "subtitle": "", "num_slides": 7
    })
    return user_prompt, "", 7


# ============================================================
# STEP 2: Generate detailed slide content from the topic
# ============================================================
def step2_generate_content(topic, subtitle, num_slides):
    """Ask the LLM to generate structured slide content."""
    emit_progress(2, f"Generating content for {num_slides} slides...")

    prompt = f"""Create detailed presentation content for the topic: "{topic}"
Subtitle: "{subtitle}"
Number of slides needed: {num_slides}

Return ONLY a JSON object with this exact structure:
{{
  "title": "{topic}",
  "subtitle": "{subtitle}",
  "slides": [
    {{"type": "section_header", "title": "Section Name"}},
    {{"type": "title_content", "title": "Slide Title", "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3", "Bullet point 4"]}},
    {{"type": "big_idea", "content": "A key takeaway or quote"}},
    {{"type": "title_content", "title": "Another Slide", "content": ["Point 1", "Point 2", "Point 3"]}}
  ]
}}

Slide types available:
- "section_header": A divider slide with just a title
- "title_content": A slide with a title and bullet points (3-5 bullets each)
- "big_idea": A slide with one big statement or key takeaway

Make the content informative, professional, and detailed. Each bullet should be a complete sentence.
Return ONLY the JSON, no extra text:"""

    raw = call_ollama(prompt)

    # Clean and parse
    raw = re.sub(r"```json\s*", "", raw)
    raw = re.sub(r"```\s*", "", raw)
    match = re.search(r"\{.*\}", raw, re.DOTALL)

    if match:
        try:
            content = json.loads(match.group())
            slides = content.get("slides", [])
            emit_progress(2, f"Content generated: {len(slides)} slides created", {
                "slide_count": len(slides)
            })
            return content
        except json.JSONDecodeError:
            pass

    # Fallback content
    emit_progress(2, "Using fallback content structure")
    return create_fallback_content(topic, subtitle)


def create_fallback_content(topic, subtitle):
    """Fallback content when LLM fails to produce valid JSON."""
    return {
        "title": topic,
        "subtitle": subtitle or "Generated Presentation",
        "slides": [
            {"type": "section_header", "title": "Overview"},
            {
                "type": "title_content",
                "title": "Introduction",
                "content": [
                    f"What is {topic}?",
                    "Historical background and evolution",
                    "Why this topic matters today",
                    "Key terminology and definitions"
                ]
            },
            {"type": "big_idea", "content": f"Understanding {topic} is essential for the future"},
            {
                "type": "title_content",
                "title": "Key Concepts",
                "content": [
                    "Core principles and foundations",
                    "Current state of the field",
                    "Major challenges and opportunities",
                    "Recent developments and breakthroughs"
                ]
            },
            {"type": "section_header", "title": "Looking Ahead"},
            {
                "type": "title_content",
                "title": "Future Directions",
                "content": [
                    "Emerging trends and innovations",
                    "Potential impact on industry and society",
                    "Areas for further exploration",
                    "Recommended resources for learning more"
                ]
            }
        ]
    }


# ============================================================
# STEP 3: Build the PPTX file from the structured content
# ============================================================
def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_styled_textbox(slide, left, top, width, height, text, font_size=18,
                        bold=False, color=COLORS['text'], align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = text
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
    return textbox


def create_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, COLORS['primary'])
    add_styled_textbox(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.5),
                       data["title"], font_size=54, bold=True,
                       color=COLORS['white'], align=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        add_styled_textbox(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.8),
                           data["subtitle"], font_size=24,
                           color=COLORS['light'], align=PP_ALIGN.CENTER)


def create_content_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, COLORS['white'])
    add_styled_textbox(slide, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
                       slide_data["title"], font_size=36, bold=True,
                       color=COLORS['dark'])
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    content = slide_data.get("content", slide_data.get("bullets", []))
    y = 2.0
    for bullet in content:
        if y > 5.0:
            break
        dot = slide.shapes.add_shape(9, Inches(1.0), Inches(y + 0.05), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLORS['primary']
        dot.line.fill.background()
        add_styled_textbox(slide, Inches(1.3), Inches(y), Inches(7.5), Inches(0.5),
                           str(bullet), font_size=18, color=COLORS['text'])
        y += 0.7


def create_big_idea_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, COLORS['light'])
    content = slide_data.get("content", slide_data.get("title", ""))
    add_styled_textbox(slide, Inches(1.5), Inches(2.0), Inches(7.0), Inches(2.5),
                       content, font_size=44, bold=True,
                       color=COLORS['primary'], align=PP_ALIGN.CENTER)


def create_section_header_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, COLORS['dark'])
    add_styled_textbox(slide, Inches(1), Inches(3), Inches(8), Inches(1.5),
                       slide_data["title"], font_size=48, bold=True,
                       color=COLORS['white'], align=PP_ALIGN.CENTER)


def step3_build_pptx(content, filename):
    """Build the PPTX file from structured content."""
    emit_progress(3, "Building presentation file...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    create_title_slide(prs, content)

    for slide_data in content.get("slides", []):
        stype = slide_data.get("type", "title_content")
        if stype == "title_content":
            create_content_slide(prs, slide_data)
        elif stype == "big_idea":
            create_big_idea_slide(prs, slide_data)
        elif stype == "section_header":
            create_section_header_slide(prs, slide_data)
        else:
            create_content_slide(prs, slide_data)

    prs.save(filename)
    slide_count = len(prs.slides)
    emit_progress(3, f"Presentation saved with {slide_count} slides", {
        "slide_count": slide_count
    })
    return slide_count


# ============================================================
# Main: orchestrate the 3 steps
# ============================================================
if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(json.dumps({"step": "error", "message": "Usage: python presentation_generator.py <prompt> <output_file>"}), flush=True)
        sys.exit(1)

    user_prompt = sys.argv[1]
    output_file = sys.argv[2]

    try:
        # Step 1: Understand the prompt
        topic, subtitle, num_slides = step1_extract_topic(user_prompt)

        # Step 2: Generate slide content
        content = step2_generate_content(topic, subtitle, num_slides)

        # Step 3: Build the PPTX
        slide_count = step3_build_pptx(content, output_file)

        # Final result
        emit_progress("done", "Presentation is ready!", {
            "success": True,
            "file": output_file,
            "slides": slide_count,
            "title": content.get("title", topic)
        })

    except Exception as e:
        emit_progress("error", str(e), {"success": False, "error": str(e)})
        sys.exit(1)
